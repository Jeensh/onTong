"""File upload/download API endpoints (binary files, images)."""

from __future__ import annotations

import hashlib
import json
import logging
import uuid
from datetime import datetime, timezone
from pathlib import Path

from fastapi import APIRouter, Depends, HTTPException, UploadFile
from fastapi.responses import FileResponse
from pydantic import BaseModel

from backend.core.config import settings

from backend.core.auth import User, get_current_user
from backend.core.auth.permission import require_admin

logger = logging.getLogger(__name__)

router = APIRouter(prefix="/api/files", tags=["files"], dependencies=[Depends(get_current_user)])

WIKI_DIR = Path(settings.wiki_dir)
ASSETS_DIR = WIKI_DIR / "assets"

ALLOWED_IMAGE_TYPES = {"image/png", "image/jpeg", "image/gif", "image/webp", "image/svg+xml"}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".svg"}
MAX_UPLOAD_SIZE = 10 * 1024 * 1024  # 10MB

# ── Image Registry injection ────────────────────────────────────────
_image_registry = None


def set_image_registry(registry) -> None:
    """Called from main.py to inject the ImageRegistry singleton."""
    global _image_registry
    _image_registry = registry


def get_image_registry():
    return _image_registry


# ── Unused image cleanup ────────────────────────────────────────────

def _find_unused_images() -> list[dict]:
    """Scan all markdown files for image references and return unreferenced assets."""
    if not ASSETS_DIR.exists():
        return []

    # Collect all asset filenames
    all_assets: dict[str, Path] = {}
    for f in ASSETS_DIR.iterdir():
        if f.is_file() and f.suffix.lower() in IMAGE_EXTENSIONS:
            all_assets[f.name] = f

    if not all_assets:
        return []

    # Scan all markdown files for references to asset filenames
    referenced: set[str] = set()
    for md_file in WIKI_DIR.rglob("*.md"):
        try:
            content = md_file.read_text(encoding="utf-8")
        except Exception:
            continue
        for name in all_assets:
            if name in content:
                referenced.add(name)

    # Return unreferenced assets
    unused = []
    for name, path in sorted(all_assets.items()):
        if name not in referenced:
            stat = path.stat()
            unused.append({
                "filename": name,
                "path": f"assets/{name}",
                "size": stat.st_size,
            })
    return unused


@router.get("/assets/unused", tags=["assets"])
async def list_unused_images(user: User = Depends(require_admin)):
    """List image files in assets/ that are not referenced by any markdown file."""
    unused = _find_unused_images()
    return {"unused": unused, "count": len(unused)}


@router.delete("/assets/unused", tags=["assets"])
async def delete_unused_images(user: User = Depends(require_admin)):
    """Delete all image files in assets/ that are not referenced by any markdown file."""
    unused = _find_unused_images()
    deleted = []
    for item in unused:
        full = ASSETS_DIR / item["filename"]
        try:
            full.unlink()
            deleted.append(item["path"])
            logger.info(f"Deleted unused image: {item['path']}")
        except Exception as e:
            logger.warning(f"Failed to delete {item['path']}: {e}")
    return {"deleted": deleted, "count": len(deleted)}


# ── Admin Image Management ──────────────────────────────────────────

@router.get("/assets/stats", tags=["assets"])
async def get_asset_stats(user: User = Depends(require_admin)):
    """Get image asset statistics (admin only)."""
    if not _image_registry:
        raise HTTPException(status_code=503, detail="Image registry not initialized")
    return _image_registry.stats()


@router.get("/assets", tags=["assets"])
async def list_assets(
    page: int = 1,
    size: int = 50,
    filter: str = "all",
    search: str = "",
    user: User = Depends(require_admin),
):
    """List image assets with pagination, filtering, search (admin only)."""
    if not _image_registry:
        raise HTTPException(status_code=503, detail="Image registry not initialized")
    if size > 100:
        size = 100

    result = _image_registry.list_entries(page=page, size=size, filter=filter, search=search)

    # Enrich items with derivatives + OCR info
    items = []
    for item_dict in result["items"]:
        item_dict["derivatives"] = _image_registry.get_derivatives_of(item_dict["filename"])
        item_dict["has_ocr"] = _has_ocr(item_dict["filename"])
        items.append(item_dict)

    return {
        "items": items,
        "total": result["total"],
        "page": result["page"],
        "pages": result["pages"],
    }


def _has_ocr(filename: str) -> bool:
    """Check if an image has OCR text in its sidecar."""
    sidecar = ASSETS_DIR / f"{filename}.meta.json"
    if not sidecar.exists():
        return False
    try:
        data = json.loads(sidecar.read_text(encoding="utf-8"))
        return bool(data.get("ocr_text") or data.get("description"))
    except Exception:
        return False


@router.delete("/assets/{filename}", tags=["assets"])
async def delete_asset(filename: str, user: User = Depends(require_admin)):
    """Delete a single image asset. Returns 409 if still referenced by documents."""
    if not _image_registry:
        raise HTTPException(status_code=503, detail="Image registry not initialized")

    entry = _image_registry.get_by_filename(filename)
    if not entry:
        raise HTTPException(status_code=404, detail=f"Image not found: {filename}")
    if entry.ref_count > 0:
        raise HTTPException(
            status_code=409,
            detail=f"Image still referenced by {entry.ref_count} document(s): {sorted(entry.referenced_by)}",
        )

    image_path = ASSETS_DIR / filename
    if image_path.exists():
        image_path.unlink()
    sidecar_path = ASSETS_DIR / f"{filename}.meta.json"
    if sidecar_path.exists():
        sidecar_path.unlink()

    freed = entry.size_bytes
    _image_registry.remove(filename)
    logger.info(f"Deleted image asset: {filename} ({freed} bytes)")
    return {"deleted": filename, "freed_bytes": freed}


@router.post("/assets/bulk-delete", tags=["assets"])
async def bulk_delete_assets(user: User = Depends(require_admin)):
    """Delete ALL unreferenced (ref_count==0) images. Returns count + freed space."""
    if not _image_registry:
        raise HTTPException(status_code=503, detail="Image registry not initialized")

    unused = _image_registry.get_unused_filenames()
    deleted_count = 0
    freed_bytes = 0

    for filename in unused:
        entry = _image_registry.get_by_filename(filename)
        if not entry:
            continue

        image_path = ASSETS_DIR / filename
        if image_path.exists():
            image_path.unlink()
        sidecar_path = ASSETS_DIR / f"{filename}.meta.json"
        if sidecar_path.exists():
            sidecar_path.unlink()

        freed_bytes += entry.size_bytes
        _image_registry.remove(filename)
        deleted_count += 1

    logger.info(f"Bulk deleted {deleted_count} unused images ({freed_bytes} bytes freed)")
    return {"deleted": deleted_count, "freed_bytes": freed_bytes}


class InheritOCRRequest(BaseModel):
    source_filename: str


@router.post("/assets/{filename}/inherit-ocr", tags=["assets"])
async def inherit_ocr(filename: str, req: InheritOCRRequest):
    """Copy OCR text + description from source image sidecar to target. Sets source field."""
    target_path = ASSETS_DIR / filename
    if not target_path.exists():
        raise HTTPException(status_code=404, detail=f"Target image not found: {filename}")

    source_sidecar = ASSETS_DIR / f"{req.source_filename}.meta.json"
    if not source_sidecar.exists():
        raise HTTPException(status_code=404, detail=f"Source sidecar not found: {req.source_filename}")

    try:
        source_data = json.loads(source_sidecar.read_text(encoding="utf-8"))
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to read source sidecar: {e}")

    from backend.application.image.models import ImageAnalysis, save_sidecar

    analysis = ImageAnalysis(
        ocr_text=source_data.get("ocr_text", ""),
        description=source_data.get("description", ""),
        provider=source_data.get("provider", "inherited"),
        ocr_engine=source_data.get("ocr_engine", "inherited"),
        processed_at=datetime.now(timezone.utc),
        source=req.source_filename,
    )
    save_sidecar(target_path, analysis)

    if _image_registry:
        entry = _image_registry.get_by_filename(filename)
        if entry:
            entry.source = req.source_filename

    return {"filename": filename, "source": req.source_filename, "inherited": True}


# ── Upload ───────────────────────────────────────────────────────────

@router.post("/upload/image")
async def upload_image(file: UploadFile):
    """Upload an image file to wiki/assets/ with SHA-256 content-hash dedup."""
    if file.content_type not in ALLOWED_IMAGE_TYPES:
        raise HTTPException(
            status_code=400,
            detail=f"Unsupported image type: {file.content_type}. Allowed: {', '.join(ALLOWED_IMAGE_TYPES)}",
        )

    data = await file.read()
    if len(data) > MAX_UPLOAD_SIZE:
        raise HTTPException(status_code=400, detail="File too large (max 10MB)")

    # Content-hash dedup
    sha256 = hashlib.sha256(data).hexdigest()

    if _image_registry:
        existing = _image_registry.get_by_hash(sha256)
        if existing:
            logger.info(f"Dedup hit: {file.filename} → existing {existing.filename}")
            return {
                "path": f"assets/{existing.filename}",
                "filename": existing.filename,
                "deduplicated": True,
            }

    ASSETS_DIR.mkdir(parents=True, exist_ok=True)

    ext = Path(file.filename or "image.png").suffix or ".png"
    filename = f"{sha256[:12]}{ext}"
    dest = ASSETS_DIR / filename

    # Handle rare hash-prefix collision
    if dest.exists():
        existing_sha = hashlib.sha256(dest.read_bytes()).hexdigest()
        if existing_sha == sha256:
            return {
                "path": f"assets/{filename}",
                "filename": filename,
                "deduplicated": True,
            }
        filename = f"{sha256}{ext}"
        dest = ASSETS_DIR / filename

    dest.write_bytes(data)

    # Register in image registry
    if _image_registry:
        from backend.application.image.image_registry import ImageEntry
        _image_registry.register(ImageEntry(
            filename=filename,
            sha256=sha256,
            size_bytes=len(data),
            width=0,
            height=0,
            ref_count=0,
            referenced_by=set(),
            source=None,
            created_at=datetime.now(timezone.utc),
        ))

    rel_path = f"assets/{filename}"
    return {"path": rel_path, "filename": filename, "deduplicated": False}


# ── PPTX slide data ─────────────────────────────────────────────────

def _parse_pptx_slides(file_path: Path) -> dict:
    """Extract slide data from a PPTX file using python-pptx."""
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation(str(file_path))
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    slides = []

    for slide in prs.slides:
        elements = []
        for shape in slide.shapes:
            el: dict = {
                "left": round(Emu(shape.left).inches, 3),
                "top": round(Emu(shape.top).inches, 3),
                "width": round(Emu(shape.width).inches, 3),
                "height": round(Emu(shape.height).inches, 3),
            }

            if shape.has_text_frame:
                paragraphs = []
                for para in shape.text_frame.paragraphs:
                    runs = []
                    for run in para.runs:
                        run_data: dict = {"text": run.text}
                        if run.font.bold:
                            run_data["bold"] = True
                        if run.font.italic:
                            run_data["italic"] = True
                        if run.font.size:
                            run_data["fontSize"] = round(Emu(run.font.size).pt, 1)
                        try:
                            if run.font.color and run.font.color.type is not None and run.font.color.rgb:
                                run_data["color"] = str(run.font.color.rgb)
                        except Exception:
                            pass
                        runs.append(run_data)
                    p_data: dict = {"runs": runs}
                    align = para.alignment
                    if align is not None:
                        align_map = {0: "left", 1: "center", 2: "right", 3: "justify"}
                        p_data["align"] = align_map.get(align, "left")
                    paragraphs.append(p_data)
                el["type"] = "text"
                el["paragraphs"] = paragraphs

            elif shape.shape_type == 13:  # Picture
                # Serve the embedded image via a data URL
                import base64
                img = shape.image
                blob = img.blob
                ct = img.content_type
                b64 = base64.b64encode(blob).decode("ascii")
                el["type"] = "image"
                el["src"] = f"data:{ct};base64,{b64}"

            else:
                # Fallback: if shape has text, treat as text
                if hasattr(shape, "text") and shape.text.strip():
                    el["type"] = "text"
                    el["paragraphs"] = [{"runs": [{"text": shape.text}]}]
                else:
                    continue

            elements.append(el)

        slides.append({"elements": elements})

    return {
        "slideWidth": round(Emu(slide_w).inches, 3),
        "slideHeight": round(Emu(slide_h).inches, 3),
        "totalSlides": len(slides),
        "slides": slides,
    }


@router.get("/pptx-data/{path:path}", tags=["pptx"])
async def get_pptx_data(path: str):
    """Parse a PPTX file and return slide data as JSON."""
    full = (WIKI_DIR / path).resolve()
    if not str(full).startswith(str(WIKI_DIR.resolve())):
        raise HTTPException(status_code=400, detail="Path traversal detected")
    if not full.exists():
        raise HTTPException(status_code=404, detail=f"File not found: {path}")
    if full.suffix.lower() not in (".pptx", ".ppt"):
        raise HTTPException(status_code=400, detail="Not a presentation file")
    try:
        return _parse_pptx_slides(full)
    except Exception as e:
        logger.error(f"Failed to parse PPTX {path}: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to parse PPTX: {e}")


# ── Binary file CRUD ─────────────────────────────────────────────────

BINARY_EXTENSIONS = {".xlsx", ".xls", ".pptx", ".ppt", ".pdf", ".png", ".jpg", ".jpeg", ".gif", ".webp", ".svg"}


@router.put("/{path:path}")
async def save_file(path: str, file: UploadFile):
    """Save a binary file to the wiki directory. Rejects .md files."""
    if path.endswith(".md"):
        raise HTTPException(status_code=400, detail="Use wiki API for markdown files")

    full = (WIKI_DIR / path).resolve()
    if not str(full).startswith(str(WIKI_DIR.resolve())):
        raise HTTPException(status_code=400, detail="Path traversal detected")

    data = await file.read()
    if len(data) > MAX_UPLOAD_SIZE:
        raise HTTPException(status_code=400, detail="File too large (max 10MB)")

    full.parent.mkdir(parents=True, exist_ok=True)
    full.write_bytes(data)
    return {"path": path, "size": len(data)}


@router.get("/{path:path}")
async def get_file(path: str):
    """Serve a binary file from the wiki directory."""
    full = (WIKI_DIR / path).resolve()
    if not str(full).startswith(str(WIKI_DIR.resolve())):
        raise HTTPException(status_code=400, detail="Path traversal detected")
    if not full.exists():
        raise HTTPException(status_code=404, detail=f"File not found: {path}")
    return FileResponse(full)
